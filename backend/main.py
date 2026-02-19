"""
Backend: PPT placeholder replacer using Excel data.

- Accepts uploaded PPTX template and Excel file
- Replaces placeholders across ALL slides/shapes while preserving formatting (best-effort)
- Exposes a simple Flask API endpoint to generate a filled PPTX

Install:
  pip install flask pandas openpyxl python-pptx flask-cors

Run:
  python main.py

POST:
  /generate
  FormData with 'template' (PPTX) and 'excel' (XLSX) files
"""

import os
import re
import uuid
from io import BytesIO
from pathlib import Path

import pandas as pd
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.dml.color import RGBColor

# ---------- Temp directory ----------
TEMP_DIR = Path(__file__).parent / "temp"
TEMP_DIR.mkdir(exist_ok=True)

# ---------- Placeholder detection ----------
def detect_placeholders(prs: Presentation) -> list[str]:
    """Extract all {{PLACEHOLDER}} patterns from the presentation."""
    placeholders = set()
    pattern = re.compile(r'\{\{([^}]+)\}\}')
    
    for slide in prs.slides:
        for shape in iter_all_shapes(slide):
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text:
                            matches = pattern.findall(run.text)
                            placeholders.update(matches)
            
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text:
                                    matches = pattern.findall(run.text)
                                    placeholders.update(matches)
    
    return sorted(placeholders)


# Mapping: PPT placeholder -> Excel column
MAPPING = {
    "[candidate_name]": "candidate_name",
    "[Name]": "candidate_name",
    "[positioning_blurb]": "positioning_blurb",
    "[deep_ai_experience_1]": "deep_ai_experience_1",
    "[deep_ai_experience_2]": "deep_ai_experience_2",
    "[deep_ai_experience_3]": "deep_ai_experience_3",
    "[deep_ai_experience_4]": "deep_ai_experience_4",
    "[key_experience_1]": "key_experience_1",
    "[key_experience_2]": "key_experience_2",
    "[key_experience_3]": "key_experience_3",
    "[key_experience_4]": "key_experience_4",
    "[quality_1_title]": "quality_1_title",
    "[quality_2_title]": "quality_2_title",
    "[quality_3_title]": "quality_3_title",
    "[quality_1_description]": "quality_1_description",
    "[quality_2_description]": "quality_2_description",
    "[quality_3_description]": "quality_3_description",
    "[backend_1]": "backend_1",
    "[backend_2]": "backend_2",
    "[backend_3]": "backend_3",
    "[backend_4]": "backend_4",
    "[backend_5]": "backend_5",
    "[backend_6]": "backend_6",
    "[full-stack_1]": "full-stack_1",
    "[full-stack_2]": "full-stack_2",
    "[full-stack_3]": "full-stack_3",
    "[full-stack_4]": "full-stack_4",
    "[full-stack_5]": "full-stack_5",
    "[full-stack_6]": "full-stack_6",
    "[ai_1]": "ai_1",
    "[ai_2]": "ai_2",
    "[ai_3]": "ai_3",
    "[ai_4]": "ai_4",
    "[ai_5]": "ai_5",
    "[ai_6]": "ai_6",
    "[project_1_name]": "project_1_name",
    "[project_1_context]": "project_1_context",
    "[project_2_name]": "project_2_name",
    "[project_2_context]": "project_2_context",
    "[project_3_name]": "project_3_name",
    "[project_3_context]": "project_3_context",
    "[project_4_name]": "project_4_name",
    "[project_4_context]": "project_4_context",
    "[project_impact_1]": "project_1_business_impact",
    "[project_1_business_impact]": "project_1_business_impact",
    "[project_impact_2]": "project_2_business_impact",
    "[project_2_business_impact]": "project_2_business_impact",
    "[Project 1]": "project_1_name",
    "[Description_p1]": "project_1_context",
    "[Project 2]": "project_2_name",
    "[Description_p2]": "project_2_context",
    "[industry_1]": "industry_1",
    "[industry_1_impact]": "industry_1_impact",
    "[industry_2]": "industry_2",
    "[industry_2_impact]": "industry_2_impact",
    "[industry_3]": "industry_3",
    "[industry_3_impact]": "industry_3_impact",
    "[industry_4]": "industry_4",
    "[industry_4_impact]": "industry_4_impact",
}

# ---------- Helpers ----------
def _normalize(s: str) -> str:
    return re.sub(r"\s+", " ", str(s or "")).strip()

def load_excel_df(excel_file) -> pd.DataFrame:
    """Load Excel from file object or path."""
    df = pd.read_excel(excel_file, engine="openpyxl")
    df.columns = [str(c).strip() for c in df.columns]
    return df

def select_row(df: pd.DataFrame, row_index: int = 0) -> pd.Series:
    """Select a row from the dataframe."""
    if row_index < 0 or row_index >= len(df):
        raise ValueError(f"row_index out of range. Got {row_index}, Excel has {len(df)} rows.")
    return df.iloc[row_index]

def build_replacements(row: pd.Series) -> dict[str, str]:
    """Build replacement dict using the mapping.

    - Normal values are normalized.
    - Empty/missing values are replaced with a red warning text like:
      'No quality_1_title field in excel'.
    """
    replacements: dict[str, str] = {}
    
    for placeholder, col in MAPPING.items():
        field_name = placeholder.strip("[]")
        if col in row.index:
            val = row[col]
            if pd.isna(val) or str(val).strip() == "":
                val = f"No {field_name} field in excel"
            else:
                val = _normalize(val)
            replacements[placeholder] = val
        else:
            replacements[placeholder] = f"No {field_name} field in excel"
    
    return replacements

from pptx.util import Pt

NAME_PLACEHOLDERS = {"[candidate_name]", "[Name]"}
NAME_COLOR = RGBColor(0, 102, 204)  # a pleasant blue for names
MISSING_FIELD_COLOR = RGBColor(220, 53, 69)  # a nice red for missing fields


def auto_fit_text(text_frame, max_font_size=None, min_font_size=8):
    """Auto-fit text by reducing font size if content overflows while preserving style."""
    if not text_frame or not text_frame.text.strip():
        return
    
    text_frame.word_wrap = True
    
    # Get original font size
    original_size = None
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                original_size = run.font.size
                break
        if original_size:
            break
    
    if not original_size:
        return
    
    # Apply max font size limit
    if max_font_size and original_size > Pt(max_font_size):
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size:
                    run.font.size = Pt(max_font_size)

def _is_missing_field_text(text: str) -> bool:
    return text.startswith("No ") and text.endswith(" field in excel")


def _copy_run_format(src_run, dst_run):
    """Copy font formatting from one run to another."""
    dst_font = dst_run.font
    src_font = src_run.font
    dst_font.name = src_font.name
    dst_font.size = src_font.size
    dst_font.bold = src_font.bold
    dst_font.italic = src_font.italic
    dst_font.underline = src_font.underline
    # src_font.color may be a special "_NoneColor" without an .rgb attribute
    try:
        rgb = getattr(src_font.color, "rgb", None)
    except Exception:
        rgb = None
    if rgb is not None:
        dst_font.color.rgb = rgb


def _split_run_for_name(paragraph, run, placeholder, replacement_text):
    """
    Replace a name placeholder inside a run so that only the name
    substring is colored, not the whole line.
    """
    original = run.text or ""
    idx = original.find(placeholder)
    if idx == -1:
        return False

    before = original[:idx]
    after = original[idx + len(placeholder) :]

    # We will transform the existing run to "before",
    # then insert name + after runs right after it.
    run.text = before

    # Name run
    name_run = paragraph.add_run()
    name_run.text = replacement_text
    _copy_run_format(run, name_run)
    name_run.font.color.rgb = NAME_COLOR
    # Insert new runs directly after the original run using underlying XML
    run._r.addnext(name_run._r)

    # After run
    if after:
        after_run = paragraph.add_run()
        after_run.text = after
        _copy_run_format(run, after_run)
        name_run._r.addnext(after_run._r)

    return True


def _highlight_name_in_paragraph(paragraph, name_text: str):
    """Color only the given name_text substring inside the paragraph."""
    if not name_text:
        return

    for run in list(paragraph.runs):
        original = run.text or ""
        idx = original.find(name_text)
        if idx == -1:
            continue

        before = original[:idx]
        after = original[idx + len(name_text) :]

        run.text = before

        # Name run
        name_run = paragraph.add_run()
        name_run.text = name_text
        _copy_run_format(run, name_run)
        name_run.font.color.rgb = NAME_COLOR
        run._r.addnext(name_run._r)

        # After run
        if after:
            after_run = paragraph.add_run()
            after_run.text = after
            _copy_run_format(run, after_run)
            name_run._r.addnext(after_run._r)

        break


def replace_in_text_runs(text_frame, replacements: dict[str, str]) -> int:
    """Replace placeholders while preserving formatting (font, color, style).

    Additionally:
    - Candidate name placeholders (`[candidate_name]`, `[Name]`) are colored `NAME_COLOR`.
    - Missing/empty fields (text like 'No quality_1_title field in excel') are colored `MISSING_FIELD_COLOR`.
    """
    count = 0
    if not text_frame:
        return 0

    placeholders = list(replacements.keys())
    if not placeholders:
        return 0

    for paragraph in text_frame.paragraphs:
        # First pass: run-level replace (preserves formatting)
        for run in paragraph.runs:
            original = run.text
            if not original:
                continue
            new_text = original
            missing_highlight = False
            for ph in placeholders:
                if ph in new_text:
                    replacement_text = replacements[ph]
                    # For name placeholders, split the run so only the
                    # name substring is colored, not the whole line.
                    if ph in NAME_PLACEHOLDERS and not _is_missing_field_text(replacement_text):
                        if _split_run_for_name(paragraph, run, ph, replacement_text):
                            new_text = run.text  # "before" part; name+after handled by new runs
                            continue
                    new_text = new_text.replace(ph, replacement_text)
                    if _is_missing_field_text(replacement_text):
                        missing_highlight = True
            if new_text != original:
                run.text = new_text
                if missing_highlight:
                    run.font.color.rgb = MISSING_FIELD_COLOR
                count += 1

        # Second pass: handle placeholders spanning runs
        para_text = "".join([r.text or "" for r in paragraph.runs])
        if any(ph in para_text for ph in placeholders):
            merged = para_text
            missing_highlight = False
            name_texts: list[str] = []
            for ph in placeholders:
                if ph in merged:
                    replacement_text = replacements[ph]
                    merged = merged.replace(ph, replacement_text)
                    if _is_missing_field_text(replacement_text):
                        missing_highlight = True
                    if ph in NAME_PLACEHOLDERS and not _is_missing_field_text(replacement_text):
                        name_texts.append(replacement_text)

            if paragraph.runs:
                # Preserve first run's formatting
                first_run = paragraph.runs[0]
                first_run.text = merged
                if missing_highlight:
                    first_run.font.color.rgb = MISSING_FIELD_COLOR
                for r in paragraph.runs[1:]:
                    r.text = ""
            else:
                new_run = paragraph.add_run()
                new_run.text = merged
                if missing_highlight:
                    new_run.font.color.rgb = MISSING_FIELD_COLOR

            # After merging, highlight each name substring individually
            for name_text in name_texts:
                _highlight_name_in_paragraph(paragraph, name_text)

            count += 1

    return count

def iter_all_shapes(slide):
    """
    Recursively iterate shapes, including grouped shapes.
    """
    for shape in slide.shapes:
        yield from _iter_shape(shape)

def _iter_shape(shape):
    yield shape
    if hasattr(shape, "shapes"):  # group shape
        for s in shape.shapes:
            yield from _iter_shape(s)

def apply_replacements_to_ppt(template_file, replacements: dict[str, str]) -> BytesIO:
    """Apply replacements to the presentation while preserving formatting."""
    prs = Presentation(template_file)

    for slide_idx, slide in enumerate(prs.slides):
        for shape in iter_all_shapes(slide):
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                replace_in_text_runs(shape.text_frame, replacements)
                
                # Auto-fit for first slide
                if slide_idx == 0:
                    auto_fit_text(shape.text_frame, max_font_size=44)
                else:
                    auto_fit_text(shape.text_frame)

            # Tables
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        replace_in_text_runs(cell.text_frame, replacements)
                        auto_fit_text(cell.text_frame)

    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out


# ---------- Flask App ----------
app = Flask(__name__)
CORS(app)

@app.get("/health")
def health():
    return jsonify({"ok": True})

@app.post("/generate")
def generate():
    """
    POST FormData:
      - template: PPTX file with {{PLACEHOLDER}} tags
      - excel: XLSX file with data
      - row_index: (optional) which row to use, defaults to 0

    Response:
      returns the filled pptx file
    """
    try:
        # Get uploaded files
        if 'template' not in request.files:
            return jsonify({"error": "Missing 'template' file"}), 400
        if 'excel' not in request.files:
            return jsonify({"error": "Missing 'excel' file"}), 400
        
        template_file = request.files['template']
        excel_file = request.files['excel']
        
        if template_file.filename == '':
            return jsonify({"error": "Empty template filename"}), 400
        if excel_file.filename == '':
            return jsonify({"error": "Empty excel filename"}), 400
        
        # Get optional row index
        row_index_str = request.form.get("row_index", "0")
        try:
            row_index = int(row_index_str)
        except (TypeError, ValueError):
            return jsonify({"error": "Invalid 'row_index'. It must be an integer."}), 400
        
        print(f"[Backend] Processing files: {template_file.filename}, {excel_file.filename}")
        print(f"[Backend] Using row index: {row_index}")
        
        # Load Excel data
        df = load_excel_df(excel_file)
        print(f"[Backend] Excel loaded: {len(df)} rows, columns: {list(df.columns)}")
        
        # Select row
        row = select_row(df, row_index=row_index)
        
        # Build replacements
        replacements = build_replacements(row)
        print(f"[Backend] Replacements: {replacements}")
        
        # Apply replacements
        template_file.seek(0)  # Reset file pointer again
        ppt_bytes = apply_replacements_to_ppt(template_file, replacements)
        
        # Save to temp directory
        output_filename = f"{uuid.uuid4()}_output.pptx"
        output_path = TEMP_DIR / output_filename
        with open(output_path, "wb") as f:
            f.write(ppt_bytes.getbuffer())
        print(f"[Backend] Saved output to: {output_path}")
        
        ppt_bytes.seek(0)
        return send_file(
            ppt_bytes,
            as_attachment=True,
            download_name="generated.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    except Exception as e:
        print(f"[Backend] Error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    print(f"[Backend] Starting server on port 8000")
    print(f"[Backend] Temp directory: {TEMP_DIR}")
    app.run(host="0.0.0.0", port=8000, debug=True)
